"""Data-grounding verification for the Software & SaaS operating memo
deck builder.

Same discipline as tests/test_build_credit_memo.py: independently
re-recalculate the source workbook here (not by re-reading the builder's
internal variables) and assert every headline figure on the deck matches
it exactly.
"""
from __future__ import annotations

import shutil
import sys
import tempfile
import unittest
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))
sys.path.insert(0, str(ROOT / "tools"))
sys.path.insert(0, str(ROOT / "tools" / "builders"))

from pptx import Presentation  # noqa: E402

from recalc import recalc  # noqa: E402
import openpyxl  # noqa: E402

from build_software_memo import build_software_memo, _load_case, _manifest  # noqa: E402


def _slide_text(prs: Presentation) -> list[str]:
    texts = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" ".join(cell.text for cell in row.cells))
        texts.append("\n".join(parts))
    return texts


class SoftwareMemoDataGroundingTests(unittest.TestCase):
    CASE_ID = "software-public-uipath-fy2023-stress"

    @classmethod
    def setUpClass(cls) -> None:
        cls.output = Path(tempfile.mkdtemp()) / "software-memo-test.pptx"
        build_software_memo(cls.CASE_ID, cls.output)
        cls.prs = Presentation(str(cls.output))
        cls.slide_text = _slide_text(cls.prs)
        cls.full_text = "\n".join(cls.slide_text)

        case = _load_case(cls.CASE_ID)
        src = ROOT / case["output"]
        scratch_dir = tempfile.mkdtemp()
        scratch = Path(scratch_dir) / src.name
        shutil.copyfile(src, scratch)
        result = recalc(str(scratch), timeout=90)
        assert result.get("status") == "success" and not result.get("total_errors", 1), result
        cls.wb = openpyxl.load_workbook(scratch, data_only=True)
        cls.manifest = _manifest(cls.CASE_ID)

    def test_deck_has_expected_slide_count(self) -> None:
        self.assertEqual(len(self.prs.slides), 6)

    def test_real_arr_figures_match_workbook_and_disclosed_ties(self) -> None:
        arr = self.wb["ARR Rollforward"]
        beginning = arr["C5"].value
        ending = arr["C10"].value
        nrr = arr["C12"].value
        # These must tie to UiPath's own disclosed figures closely (not
        # just "be a number") -- the whole point of using real ARR. Ending
        # ARR carries a known ~$0.11mm rounding artifact from the case's
        # own 4-decimal-rounded derived flow rates (New/Expansion), not a
        # memo-builder bug -- delta reflects that, not loosened blindly.
        self.assertAlmostEqual(beginning, 925.3, places=1)
        self.assertAlmostEqual(ending, 1203.8, delta=0.15)
        self.assertAlmostEqual(nrr, 1.23, places=2)
        self.assertIn(f"${beginning:,.1f}mm", self.slide_text[2])
        self.assertIn(f"${ending:,.1f}mm", self.slide_text[2])

    def test_realized_outcome_matches_manifest_exactly(self) -> None:
        realized = self.manifest["outcome"]["realized"]
        self.assertEqual(realized, 1.19)
        self.assertIn(f"{realized * 100:.0f}%", self.slide_text[5])

    def test_real_vs_illustrative_labels_present(self) -> None:
        self.assertIn("Illustrative", self.full_text)
        self.assertIn("Derived", self.full_text)

    def test_overall_checks_status_matches_workbook(self) -> None:
        checks = self.wb["Checks"]
        overall = checks["C18"].value
        self.assertIn(f"Overall checks: {overall}", self.slide_text[5])

    @classmethod
    def tearDownClass(cls) -> None:
        shutil.rmtree(cls.output.parent, ignore_errors=True)


if __name__ == "__main__":
    unittest.main()
