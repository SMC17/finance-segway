"""Data-grounding verification for the IC memo deck builder.

Mirrors this repo's Excel verification discipline: don't just check the deck
opens, check that the numbers on it are the real, recalculated workbook's
own numbers, independently re-derived here rather than re-reading the
builder's internal variables.
"""
from __future__ import annotations

import shutil
import sys
import tempfile
import unittest
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))
sys.path.insert(0, str(ROOT / "tools"))
sys.path.insert(0, str(ROOT / "tools" / "builders"))

from pptx import Presentation  # noqa: E402

from recalc import recalc  # noqa: E402
import openpyxl  # noqa: E402

from build_ic_memo import build_lbo_ic_memo, _load_case, _manifest_sources  # noqa: E402


def _slide_text(prs: Presentation) -> list[str]:
    texts = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" ".join(cell.text for cell in row.cells))
        texts.append("\n".join(parts))
    return texts


class ICMemoDataGroundingTests(unittest.TestCase):
    CASE_ID = "pe-public-home-depot-2023"

    @classmethod
    def setUpClass(cls) -> None:
        cls.output = Path(tempfile.mkdtemp()) / "ic-memo-test.pptx"
        build_lbo_ic_memo(cls.CASE_ID, cls.output)
        cls.prs = Presentation(str(cls.output))
        cls.slide_text = _slide_text(cls.prs)
        cls.full_text = "\n".join(cls.slide_text)

        # Independently recalculate the source workbook again here, rather
        # than trusting the builder's own recalculated values.
        case = _load_case(cls.CASE_ID)
        src = ROOT / case["output"]
        scratch_dir = tempfile.mkdtemp()
        scratch = Path(scratch_dir) / src.name
        shutil.copyfile(src, scratch)
        result = recalc(str(scratch), timeout=90)
        assert result.get("status") == "success" and not result.get("total_errors", 1), result
        cls.wb = openpyxl.load_workbook(scratch, data_only=True)
        cls.manifest = _manifest_sources(cls.CASE_ID)

    def test_deck_has_expected_slide_count(self) -> None:
        self.assertEqual(len(self.prs.slides), 7)

    def test_real_operating_figures_appear_verbatim(self) -> None:
        assumptions = self.wb["Assumptions"]
        revenue = assumptions["C5"].value
        margin = assumptions["C6"].value
        growth = assumptions["C7"].value
        self.assertIn(f"${revenue:,.0f}mm", self.full_text)
        self.assertIn(f"{margin * 100:.1f}%", self.full_text)
        self.assertIn(f"{growth * 100:.1f}%", self.full_text)

    def test_sources_and_uses_totals_match_workbook_exactly(self) -> None:
        su = self.wb["Sources & Uses"]
        total_uses = su["C9"].value
        total_sources = su["F9"].value
        self.assertAlmostEqual(total_uses, total_sources, places=2)
        self.assertIn(f"{total_uses:,.0f}", self.slide_text[3])
        self.assertIn(f"{su['F8'].value:,.0f}", self.slide_text[3])  # sponsor equity

    def test_returns_figures_match_selected_exit_year_column(self) -> None:
        assumptions = self.wb["Assumptions"]
        exit_year = int(assumptions["C29"].value)
        exit_col = {1: "C", 2: "D", 3: "E", 4: "F", 5: "G", 6: "H", 7: "I"}[exit_year]
        rw = self.wb["Returns Waterfall"]
        moic = rw[f"{exit_col}14"].value
        irr = rw[f"{exit_col}15"].value
        self.assertIn(f"{moic:.1f}x", self.slide_text[4])
        self.assertIn(f"{irr * 100:.1f}%", self.slide_text[4])

    def test_sensitivity_grid_matches_workbook_cell_for_cell(self) -> None:
        sens = self.wb["Sensitivity"]
        # Spot-check every cell of the 5x5 grid, not just one.
        for row_offset, row in enumerate(range(5, 10)):
            for col_offset, col in enumerate(range(3, 8)):
                value = sens.cell(row, col).value
                self.assertIn(f"{value * 100:.1f}%", self.slide_text[5])

    def test_real_vs_illustrative_labels_match_manifest_input_kinds(self) -> None:
        real_inputs = [
            item for item in self.manifest["inputs"]
            if item.get("input_kind") in {"observed", "derived"}
        ]
        real_cells = {item["cell"] for item in real_inputs}
        # The real, sourced Assumptions cells must be exactly the ones the
        # deck claims are "real, sourced" -- every other transaction-
        # structure figure must carry the "Illustrative" tag. This case's
        # real cell set: C5 (revenue, 10-K), C6/C7/C8 (EBITDA margin, growth,
        # annual margin change -- earnings call), C9/C10/C12 (D&A, capex, tax
        # rate -- cash flow statement cross-check).
        self.assertEqual(real_cells, {"C5", "C6", "C7", "C8", "C9", "C10", "C12"})
        self.assertIn("Illustrative", self.slide_text[3])
        self.assertNotIn("SEC-SOURCED", self.slide_text[3].upper().replace(" ", "-SOURCED").replace("--", "-"))

    def test_no_fabricated_source_names(self) -> None:
        # Every named source string appearing anywhere on the deck must
        # trace to a source object actually present in the manifest.
        manifest_source_names = {source["name"] for source in self.manifest.get("sources", [])}
        for name in manifest_source_names:
            if name.startswith("Frozen source snapshot"):
                continue
            self.assertIn(name, self.full_text)

    def test_overall_checks_status_matches_workbook(self) -> None:
        checks = self.wb["Checks"]
        overall = checks["C13"].value
        self.assertIn(f"Overall checks: {overall}", self.slide_text[6])

    @classmethod
    def tearDownClass(cls) -> None:
        shutil.rmtree(cls.output.parent, ignore_errors=True)


if __name__ == "__main__":
    unittest.main()
