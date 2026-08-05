"""Data-grounding verification for the Private Credit memo deck builder.

Same discipline as tests/test_build_ic_memo.py: independently re-recalculate
the source workbook here (not by re-reading the builder's internal
variables) and assert every headline figure on the deck matches it exactly.
"""
from __future__ import annotations

import re
import shutil
import sys
import tempfile
import unittest
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))
sys.path.insert(0, str(ROOT / "tools"))
sys.path.insert(0, str(ROOT / "tools" / "builders"))

from pptx import Presentation  # noqa: E402

from recalc import recalc  # noqa: E402
import openpyxl  # noqa: E402

from build_credit_memo import build_credit_memo, _load_case, _manifest_sources  # noqa: E402


def _slide_text(prs: Presentation) -> list[str]:
    texts = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" ".join(cell.text for cell in row.cells))
        texts.append("\n".join(parts))
    return texts


class CreditMemoDataGroundingTests(unittest.TestCase):
    CASE_ID = "credit-public-ares-2024"

    @classmethod
    def setUpClass(cls) -> None:
        cls.output = Path(tempfile.mkdtemp()) / "credit-memo-test.pptx"
        build_credit_memo(cls.CASE_ID, cls.output)
        cls.prs = Presentation(str(cls.output))
        cls.slide_text = _slide_text(cls.prs)
        cls.full_text = "\n".join(cls.slide_text)

        case = _load_case(cls.CASE_ID)
        src = ROOT / case["output"]
        scratch_dir = tempfile.mkdtemp()
        scratch = Path(scratch_dir) / src.name
        shutil.copyfile(src, scratch)
        result = recalc(str(scratch), timeout=90)
        assert result.get("status") == "success" and not result.get("total_errors", 1), result
        cls.wb = openpyxl.load_workbook(scratch, data_only=True)
        cls.manifest = _manifest_sources(cls.CASE_ID)

    def test_deck_has_expected_slide_count(self) -> None:
        self.assertEqual(len(self.prs.slides), 7)

    def test_real_portfolio_exposure_figures_appear_verbatim(self) -> None:
        portfolio = self.wb["Portfolio & Concentration"]
        for row in range(5, 10):
            exposure = portfolio.cell(row, 3).value
            self.assertIn(f"{exposure:,.1f}", self.full_text)
        total = portfolio["C12"].value
        self.assertIn(f"${total:,.0f}mm", self.full_text)

    def _chart_series(self, slide_index: int) -> dict[str, list[float]]:
        series_by_name: dict[str, list[float]] = {}
        for shape in self.prs.slides[slide_index].shapes:
            if shape.has_chart:
                for plot in shape.chart.plots:
                    for series in plot.series:
                        series_by_name[series.name] = list(series.values)
        return series_by_name

    def test_covenant_headroom_series_match_workbook_cell_for_cell(self) -> None:
        covenants = self.wb["Covenants"]
        expected_leverage = [covenants.cell(7, c).value for c in range(3, 8)]
        expected_coverage = [covenants.cell(10, c).value for c in range(3, 8)]
        expected_dscr = [covenants.cell(13, c).value for c in range(3, 8)]

        series = self._chart_series(3)
        for actual, expected in zip(series["Leverage headroom"], expected_leverage):
            self.assertAlmostEqual(actual, expected, places=6)
        for actual, expected in zip(series["Coverage headroom"], expected_coverage):
            self.assertAlmostEqual(actual, expected, places=6)
        for actual, expected in zip(series["DSCR headroom"], expected_dscr):
            self.assertAlmostEqual(actual, expected, places=6)

        # Direct, checkable figures on the slide: Year 1 stats.
        self.assertIn(f"{covenants['C5'].value:.2f}x", self.slide_text[3])
        self.assertIn(f"{covenants['C6'].value:.2f}x", self.slide_text[3])
        self.assertIn(covenants["C14"].value, self.slide_text[3])

    def test_recovery_bridge_matches_workbook_exactly(self) -> None:
        recovery = self.wb["Recovery"]
        rows = [6, 9, 10]
        for row in rows:
            base = recovery.cell(row, 3).value
            downside = recovery.cell(row, 4).value
            for value in (base, downside):
                self.assertTrue(
                    f"{value:.1f}x" in self.slide_text[5] or f"{value:,.1f}" in self.slide_text[5],
                    f"row {row} value {value} not found on recovery slide",
                )

    def test_real_vs_illustrative_labels_match_manifest_input_kinds(self) -> None:
        real_inputs = [
            item for item in self.manifest["inputs"]
            if item.get("input_kind") in {"observed", "derived"}
        ]
        real_cells = {item["cell"] for item in real_inputs}
        self.assertEqual(real_cells, {"C5", "C6", "C7", "C8", "C9"})
        self.assertIn("Illustrative", self.slide_text[3])

    def test_no_fabricated_source_names(self) -> None:
        # Guard against fabrication: every "Ares Capital ... Form 10-<K/Q>"
        # string on the deck must trace to a source object actually present
        # in the manifest. (The reverse -- every manifest source must appear
        # on the deck -- doesn't hold here: the manifest also lists the
        # second-quarter 2025 10-Q used only for the outcome/realized-value
        # check, which this deck doesn't display, so it's legitimately
        # absent from the slide text.)
        manifest_source_names = {source["name"] for source in self.manifest.get("sources", [])}
        # The builder cites sources in two fixed forms: a footer
        # "Source: <name> — <url>" and inline "from <name>." Extract those
        # exact citations rather than any substring containing "Form 10-".
        cited = set(re.findall(r"Source: ([^—\n]+?) — ", self.full_text))
        cited |= set(re.findall(r"from (Ares Capital \d{4} Form 10-K)\.", self.full_text))
        self.assertTrue(cited, "expected at least one source citation on the deck")
        for name in cited:
            self.assertIn(name.strip(), manifest_source_names)

    def test_overall_checks_status_matches_workbook(self) -> None:
        checks = self.wb["Checks"]
        overall = checks["C9"].value
        self.assertIn(f"Overall checks: {overall}", self.slide_text[6])

    @classmethod
    def tearDownClass(cls) -> None:
        shutil.rmtree(cls.output.parent, ignore_errors=True)


if __name__ == "__main__":
    unittest.main()
