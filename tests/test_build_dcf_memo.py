"""Data-grounding verification for the Investment Banking DCF memo deck
builder.

Same discipline as tests/test_build_credit_memo.py, adapted for an
agent-tool-drafted instance (no standards/public_cases/index.json entry,
no input_kind field, no outcome) rather than a registered public case.
"""
from __future__ import annotations

import shutil
import sys
import tempfile
import unittest
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))
sys.path.insert(0, str(ROOT / "tools"))
sys.path.insert(0, str(ROOT / "tools" / "builders"))

from pptx import Presentation  # noqa: E402

from recalc import recalc  # noqa: E402
import openpyxl  # noqa: E402

from build_dcf_memo import build_dcf_memo, _load_manifest  # noqa: E402


def _slide_text(prs: Presentation) -> list[str]:
    texts = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" ".join(cell.text for cell in row.cells))
        texts.append("\n".join(parts))
    return texts


class DcfMemoDataGroundingTests(unittest.TestCase):
    SLUG = "public_adobe_dcf_proxy"

    @classmethod
    def setUpClass(cls) -> None:
        cls.output = Path(tempfile.mkdtemp()) / "dcf-memo-test.pptx"
        build_dcf_memo(cls.SLUG, cls.output)
        cls.prs = Presentation(str(cls.output))
        cls.slide_text = _slide_text(cls.prs)
        cls.full_text = "\n".join(cls.slide_text)

        cls.manifest = _load_manifest(cls.SLUG)
        src = ROOT / cls.manifest["output"]
        scratch_dir = tempfile.mkdtemp()
        scratch = Path(scratch_dir) / src.name
        shutil.copyfile(src, scratch)
        result = recalc(str(scratch), timeout=90)
        assert result.get("status") == "success" and not result.get("total_errors", 1), result
        cls.wb = openpyxl.load_workbook(scratch, data_only=True)

    def test_deck_has_expected_slide_count(self) -> None:
        self.assertEqual(len(self.prs.slides), 5)

    def test_manifest_is_agent_tool_draft_not_a_registered_public_case(self) -> None:
        self.assertEqual(self.manifest["classification"], "agent_tool_draft")
        self.assertFalse(self.manifest["counts_toward_M4"])
        self.assertIn("agent_tool_draft", self.full_text)

    def test_base_revenue_seed_matches_workbook_exactly(self) -> None:
        # Regression coverage for the real bug this same fixture caught:
        # IS!E5 must actually carry the real base revenue, not the 0
        # default that silently zeroed the whole DCF chain before the fix.
        is_sheet = self.wb["IS"]
        base_revenue = is_sheet["E5"].value
        self.assertAlmostEqual(base_revenue, 23769.0, places=0)
        self.assertIn(f"${base_revenue:,.1f}mm", self.slide_text[2])

    def test_wacc_and_terminal_growth_are_on_dcf_sheet_and_labeled_precedent(self) -> None:
        dcf = self.wb["DCF"]
        self.assertEqual(dcf["I5"].value, 0.09)
        self.assertEqual(dcf["I6"].value, 0.025)
        self.assertIn("Repo precedent", self.full_text)
        self.assertIn("no issuer", self.full_text.lower())

    def test_implied_value_per_share_is_positive_and_matches_workbook(self) -> None:
        dcf = self.wb["DCF"]
        implied = dcf["I14"].value
        self.assertGreater(implied, 0)
        self.assertIn(f"${implied:,.2f}", self.slide_text[4])

    def test_comps_gap_is_explicitly_disclosed(self) -> None:
        self.assertIn("Comps", self.full_text)
        self.assertIn("gap", self.full_text.lower())

    @classmethod
    def tearDownClass(cls) -> None:
        shutil.rmtree(cls.output.parent, ignore_errors=True)


if __name__ == "__main__":
    unittest.main()
