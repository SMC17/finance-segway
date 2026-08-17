"""Data-grounding verification for the Distressed & Restructuring memo
deck builder.

Same discipline as tests/test_build_credit_memo.py, adapted for this
case's deliberate sparsity: distressed-public-hertz-2021-reorganization
sources exactly one real fact, so the tests here specifically guard
against the deck ever implying more real coverage than that.
"""
from __future__ import annotations

import shutil
import sys
import tempfile
import unittest
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))
sys.path.insert(0, str(ROOT / "tools"))
sys.path.insert(0, str(ROOT / "tools" / "builders"))

from pptx import Presentation  # noqa: E402

from recalc import recalc  # noqa: E402
import openpyxl  # noqa: E402

from build_restructuring_memo import build_restructuring_memo, _load_case, _manifest  # noqa: E402


def _slide_text(prs: Presentation) -> list[str]:
    texts = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" ".join(cell.text for cell in row.cells))
        texts.append("\n".join(parts))
    return texts


class RestructuringMemoDataGroundingTests(unittest.TestCase):
    CASE_ID = "distressed-public-hertz-2021-reorganization"

    @classmethod
    def setUpClass(cls) -> None:
        cls.output = Path(tempfile.mkdtemp()) / "restructuring-memo-test.pptx"
        build_restructuring_memo(cls.CASE_ID, cls.output)
        cls.prs = Presentation(str(cls.output))
        cls.slide_text = _slide_text(cls.prs)
        cls.full_text = "\n".join(cls.slide_text)

        case = _load_case(cls.CASE_ID)
        src = ROOT / case["output"]
        scratch_dir = tempfile.mkdtemp()
        scratch = Path(scratch_dir) / src.name
        shutil.copyfile(src, scratch)
        result = recalc(str(scratch), timeout=90)
        assert result.get("status") == "success" and not result.get("total_errors", 1), result
        cls.wb = openpyxl.load_workbook(scratch, data_only=True)
        cls.manifest = _manifest(cls.CASE_ID)

    def test_deck_has_expected_slide_count(self) -> None:
        self.assertEqual(len(self.prs.slides), 5)

    def test_exactly_one_real_input_and_deck_says_so(self) -> None:
        real_inputs = [
            item for item in self.manifest["inputs"]
            if item.get("input_kind") in {"observed", "derived"}
        ]
        self.assertEqual(len(real_inputs), 1)
        self.assertEqual(real_inputs[0]["cell"], "C5")
        self.assertIn("exactly one real fact", self.full_text)

    def test_real_commitment_figure_matches_workbook_exactly(self) -> None:
        new_money = self.wb["New Money"]
        commitment = new_money["C5"].value
        self.assertEqual(commitment, 7500.0)
        self.assertIn(f"${commitment:,.1f}mm", self.slide_text[2])
        self.assertIn(f"${commitment:,.1f}mm", self.slide_text[3])

    def test_realized_outcome_matches_manifest_exactly(self) -> None:
        self.assertEqual(self.manifest["outcome"]["realized"], 1.0)
        self.assertIn("Completed", self.slide_text[2])

    def test_deck_does_not_claim_recovery_waterfall_is_real(self) -> None:
        # Recovery Waterfall / 13-Week Liquidity / Liquidation-vs-Reorg are
        # untouched template defaults for this case -- the deck must not
        # present any figure sourced from those sheets as if it were real.
        real_cells = {item["cell"] for item in self.manifest["inputs"]}
        self.assertNotIn("Recovery Waterfall", str(self.manifest["inputs"]))
        self.assertIn("template defaults", self.full_text)

    def test_overall_checks_status_matches_workbook(self) -> None:
        checks = self.wb["Decision & Checks"]
        overall = checks["C14"].value
        self.assertIn(f"Overall checks: {overall}", self.slide_text[4])

    @classmethod
    def tearDownClass(cls) -> None:
        shutil.rmtree(cls.output.parent, ignore_errors=True)


if __name__ == "__main__":
    unittest.main()
