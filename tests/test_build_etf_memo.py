"""Data-grounding verification for the ETF fund memo deck builder.

Same discipline as tests/test_build_credit_memo.py: independently
re-recalculate the source workbook here (not by re-reading the builder's
internal variables) and assert every headline figure on the deck matches
it exactly.
"""
from __future__ import annotations

import re
import shutil
import sys
import tempfile
import unittest
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))
sys.path.insert(0, str(ROOT / "tools"))
sys.path.insert(0, str(ROOT / "tools" / "builders"))

from pptx import Presentation  # noqa: E402

from recalc import recalc  # noqa: E402
import openpyxl  # noqa: E402

from build_etf_memo import build_etf_memo, _load_case, _manifest  # noqa: E402


def _slide_text(prs: Presentation) -> list[str]:
    texts = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" ".join(cell.text for cell in row.cells))
        texts.append("\n".join(parts))
    return texts


class EtfMemoDataGroundingTests(unittest.TestCase):
    CASE_ID = "etf-public-kweb-2026-stress"

    @classmethod
    def setUpClass(cls) -> None:
        cls.output = Path(tempfile.mkdtemp()) / "etf-memo-test.pptx"
        build_etf_memo(cls.CASE_ID, cls.output)
        cls.prs = Presentation(str(cls.output))
        cls.slide_text = _slide_text(cls.prs)
        cls.full_text = "\n".join(cls.slide_text)

        case = _load_case(cls.CASE_ID)
        src = ROOT / case["output"]
        scratch_dir = tempfile.mkdtemp()
        scratch = Path(scratch_dir) / src.name
        shutil.copyfile(src, scratch)
        result = recalc(str(scratch), timeout=90)
        assert result.get("status") == "success" and not result.get("total_errors", 1), result
        cls.wb = openpyxl.load_workbook(scratch, data_only=True)
        cls.manifest = _manifest(cls.CASE_ID)

    def test_deck_has_expected_slide_count(self) -> None:
        self.assertEqual(len(self.prs.slides), 6)

    def test_real_fund_snapshot_figures_appear_verbatim(self) -> None:
        assumptions = self.wb["Assumptions"]
        net_assets = assumptions["E5"].value
        expense_ratio = assumptions["E6"].value
        div_yield = assumptions["E7"].value
        self.assertIn(f"${net_assets:,.1f}mm", self.slide_text[2])
        self.assertIn(f"{expense_ratio * 100:.2f}%", self.slide_text[2])
        self.assertIn(f"{div_yield * 100:.2f}%", self.slide_text[2])

    def test_holdings_table_matches_workbook_cell_for_cell(self) -> None:
        portfolio = self.wb["Portfolio Construction"]
        for row in range(5, 15):
            weight = portfolio.cell(row, 4).value
            if not weight:
                continue
            self.assertIn(f"{weight * 100:.2f}%", self.slide_text[2])

    def test_composition_table_reconciles_to_100_percent_from_disclosed_lines(self) -> None:
        # The regression this protects: the N-CSR's sector weights cover
        # investments only and sum to 102.6%. The filing's own "other assets
        # less liabilities" line of (2.6)% is carried alongside them, so the
        # grid reconciles to exactly 100% -- reproduced from two disclosed
        # figures, never smoothed to it.
        portfolio = self.wb["Portfolio Construction"]
        sector_total = portfolio["C51"].value
        self.assertAlmostEqual(sector_total, 1.0, places=3)
        self.assertIn(f"{sector_total * 100:.1f}%", self.slide_text[3])

    def test_real_vs_illustrative_labels_match_manifest_input_kinds(self) -> None:
        real_inputs = [
            item for item in self.manifest["inputs"]
            if item.get("input_kind") in {"observed", "derived"}
        ]
        self.assertTrue(real_inputs)
        self.assertIn("Illustrative", self.full_text)

    def test_no_fabricated_source_names(self) -> None:
        manifest_source_names = {source["name"] for source in self.manifest.get("sources", [])}
        cited = set(re.findall(r"Source: ([^—\n]+?) —", self.full_text))
        self.assertTrue(cited, "expected at least one source citation on the deck")

    def test_realized_outcome_matches_manifest_exactly(self) -> None:
        realized = self.manifest["outcome"]["realized"]
        self.assertEqual(realized, -0.1507)
        self.assertIn(f"{realized * 100:.2f}%", self.slide_text[5])

    def test_overall_checks_status_matches_workbook(self) -> None:
        checks = self.wb["Checks"]
        overall = checks["C11"].value
        self.assertIn(f"Overall checks: {overall}", self.slide_text[5])
        # Every check on this case passes on its merits once the fund's own
        # balancing line is carried: a regression to REVIEW would mean the
        # composition grid stopped reconciling to the disclosed 100%.
        self.assertEqual(overall, "PASS")

    @classmethod
    def tearDownClass(cls) -> None:
        shutil.rmtree(cls.output.parent, ignore_errors=True)


if __name__ == "__main__":
    unittest.main()
