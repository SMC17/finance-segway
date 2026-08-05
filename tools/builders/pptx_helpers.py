"""Shared styling and layout helpers for governed IC-memo deck builders.

Mirrors institutional_helpers.py's role for Excel builders: a small, reused
set of primitives (palette, title slide, stat callouts, tables, native
charts) so every deck this repo produces looks like one system instead of
each builder inventing its own layout. Every number placed on a slide by a
caller of this module is expected to come from a real, recalculated
workbook cell or a cited source -- these helpers only render, they never
invent a figure.
"""
from __future__ import annotations

from dataclasses import dataclass
from typing import Sequence

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# "Midnight Executive" palette -- one dominant color (navy), one supporting
# tone (ice blue), one sharp accent (white on dark, navy on light).
NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE_BLUE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL = RGBColor(0x22, 0x27, 0x33)
MUTED_GRAY = RGBColor(0x6B, 0x72, 0x80)
CARD_FILL = RGBColor(0xF4, 0xF6, 0xFC)
POSITIVE = RGBColor(0x1B, 0x7A, 0x43)
NEGATIVE = RGBColor(0xA3, 0x1E, 0x2B)
ILLUSTRATIVE_TAG = RGBColor(0x8A, 0x5A, 0x00)

HEADER_FONT = "Cambria"
BODY_FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def _set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_background(slide, color: RGBColor) -> None:
    left = top = Emu(0)
    rect = slide.shapes.add_shape(1, left, top, SLIDE_W, SLIDE_H)  # MSO_SHAPE.RECTANGLE
    _set_fill(rect, color)
    rect.shadow.inherit = False
    # push to back
    slide.shapes._spTree.remove(rect._element)
    slide.shapes._spTree.insert(2, rect._element)


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = 0
    box.text_frame.margin_right = 0
    box.text_frame.margin_top = 0
    box.text_frame.margin_bottom = 0
    return box


def add_text(
    slide,
    left, top, width, height,
    text: str,
    *,
    size: int = 14,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = CHARCOAL,
    font: str = BODY_FONT,
    align=PP_ALIGN.LEFT,
):
    box = _textbox(slide, left, top, width, height)
    tf = box.text_frame
    tf.text = text
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return box


def add_bullets(
    slide,
    left, top, width, height,
    items: Sequence[str],
    *,
    size: int = 14,
    color: RGBColor = CHARCOAL,
    font: str = BODY_FONT,
    space_after: int = 8,
):
    box = _textbox(slide, left, top, width, height)
    tf = box.text_frame
    for index, item in enumerate(items):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = item
        paragraph.space_after = Pt(space_after)
        run = paragraph.runs[0]
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = font
    return box


def add_title_slide(prs: Presentation, title: str, subtitle: str, footer: str) -> None:
    slide = _blank_slide(prs)
    add_background(slide, NAVY)
    add_text(
        slide, MARGIN, Inches(2.7), SLIDE_W - 2 * MARGIN, Inches(1.4),
        title, size=40, bold=True, color=WHITE, font=HEADER_FONT,
    )
    add_text(
        slide, MARGIN, Inches(4.05), SLIDE_W - 2 * MARGIN, Inches(1.0),
        subtitle, size=18, color=ICE_BLUE, font=BODY_FONT,
    )
    add_text(
        slide, MARGIN, SLIDE_H - Inches(0.9), SLIDE_W - 2 * MARGIN, Inches(0.5),
        footer, size=11, color=ICE_BLUE, font=BODY_FONT,
    )


def add_section_header(slide, title: str, kicker: str | None = None) -> None:
    if kicker:
        add_text(
            slide, MARGIN, Inches(0.45), SLIDE_W - 2 * MARGIN, Inches(0.35),
            kicker.upper(), size=11, bold=True, color=MUTED_GRAY, font=BODY_FONT,
        )
        title_top = Inches(0.78)
    else:
        title_top = Inches(0.5)
    add_text(
        slide, MARGIN, title_top, SLIDE_W - 2 * MARGIN, Inches(0.7),
        title, size=28, bold=True, color=NAVY, font=HEADER_FONT,
    )


def add_content_slide(prs: Presentation, title: str, kicker: str | None = None):
    slide = _blank_slide(prs)
    add_background(slide, WHITE)
    add_section_header(slide, title, kicker)
    return slide


@dataclass
class Stat:
    value: str
    label: str
    tag: str | None = None  # e.g. "SEC 10-K" or "Illustrative"
    color: RGBColor = NAVY


def add_stat_row(slide, top, stats: Sequence[Stat]) -> None:
    n = len(stats)
    gap = Inches(0.35)
    usable = SLIDE_W - 2 * MARGIN - gap * (n - 1)
    card_w = Emu(int(usable / n))
    for index, stat in enumerate(stats):
        left = MARGIN + index * (card_w + gap)
        card = slide.shapes.add_shape(1, left, top, card_w, Inches(1.55))
        _set_fill(card, CARD_FILL)
        card.shadow.inherit = False
        add_text(
            slide, left + Inches(0.15), top + Inches(0.18), card_w - Inches(0.3), Inches(0.6),
            stat.value, size=26, bold=True, color=stat.color, font=HEADER_FONT,
        )
        add_text(
            slide, left + Inches(0.15), top + Inches(0.82), card_w - Inches(0.3), Inches(0.4),
            stat.label, size=11, color=CHARCOAL, font=BODY_FONT,
        )
        if stat.tag:
            add_text(
                slide, left + Inches(0.15), top + Inches(1.16), card_w - Inches(0.3), Inches(0.3),
                stat.tag, size=9, italic=True,
                color=POSITIVE if "SEC" in stat.tag or "10-K" in stat.tag or "10-Q" in stat.tag else ILLUSTRATIVE_TAG,
                font=BODY_FONT,
            )


def add_table(
    slide,
    left, top, width, height,
    headers: Sequence[str],
    rows: Sequence[Sequence[str]],
    *,
    col_widths: Sequence[float] | None = None,
    header_fill: RGBColor = NAVY,
    header_color: RGBColor = WHITE,
    body_size: int = 11,
):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    graphic_frame = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = graphic_frame.table
    if col_widths:
        total = sum(col_widths)
        for index, weight in enumerate(col_widths):
            table.columns[index].width = Emu(int(width * weight / total))
    for col_index, text in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT if col_index == 0 else PP_ALIGN.RIGHT
        run = paragraph.runs[0]
        run.font.bold = True
        run.font.size = Pt(11)
        run.font.color.rgb = header_color
        run.font.name = BODY_FONT
    for row_index, row in enumerate(rows, start=1):
        for col_index, text in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = str(text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_index % 2 else CARD_FILL
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT if col_index == 0 else PP_ALIGN.RIGHT
            run = paragraph.runs[0]
            run.font.size = Pt(body_size)
            run.font.color.rgb = CHARCOAL
            run.font.name = BODY_FONT
    return graphic_frame


def add_line_chart(
    slide,
    left, top, width, height,
    title: str,
    categories: Sequence[str],
    series: Sequence[tuple[str, Sequence[float]]],
) -> None:
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series:
        chart_data.add_series(name, values)
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, left, top, width, height, chart_data
    )
    chart = graphic_frame.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
    if len(series) > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    else:
        chart.has_legend = False
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(10)
    category_axis.format.line.color.rgb = MUTED_GRAY
    value_axis = chart.value_axis
    value_axis.tick_labels.font.size = Pt(10)
    value_axis.has_major_gridlines = True
    value_axis.major_gridlines.format.line.color.rgb = CARD_FILL
    for plot in chart.plots:
        plot.has_data_labels = False


def add_footer(slide, text: str) -> None:
    add_text(
        slide, MARGIN, SLIDE_H - Inches(0.45), SLIDE_W - 2 * MARGIN, Inches(0.35),
        text, size=9, color=MUTED_GRAY, font=BODY_FONT,
    )


def add_disclosure_note(slide, text: str, *, color: RGBColor = ILLUSTRATIVE_TAG) -> None:
    """A visible, non-decorative provenance note -- not a footer/citation
    line but a callout distinguishing real sourced figures from illustrative
    modeler assumptions on the same slide."""
    box = slide.shapes.add_shape(1, MARGIN, Inches(6.55), SLIDE_W - 2 * MARGIN, Inches(0.55))
    _set_fill(box, CARD_FILL)
    box.shadow.inherit = False
    add_text(
        slide, MARGIN + Inches(0.15), Inches(6.55) + Inches(0.1), SLIDE_W - 2 * MARGIN - Inches(0.3), Inches(0.35),
        text, size=10, italic=True, color=color, font=BODY_FONT,
    )


def save(prs: Presentation, output) -> None:
    from pathlib import Path

    output = Path(output)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
